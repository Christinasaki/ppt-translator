#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT一键翻译 - 独立脚本版本
Universal PowerPoint translation tool - Standalone version

Usage:
    python ppt_translate_standalone.py input.pptx -o output.pptx --terms terms.xlsx

Dependencies:
    pip install python-pptx Pillow openpyxl lxml
"""

import argparse
import json
import os
import platform
import re
import sys
from pathlib import Path
from typing import Optional, Dict, Set, Tuple

# Third-party imports
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree
from PIL import ImageFont
import openpyxl


# ============================================================================
# Constants
# ============================================================================

CN_RE = re.compile(r"[\u4e00-\u9fff]")
ALPHA_RE = re.compile(r"[A-Za-z0-9]+")

DEFAULT_SMALL_WORDS = {
    "a", "an", "the", "and", "but", "or", "nor", "for", "so", "yet",
    "at", "by", "from", "in", "into", "of", "off", "on", "onto", "per",
    "to", "up", "via", "with", "as"
}

DEFAULT_PRESERVED_TOKENS = {
    "ai": "AI", "api": "API", "app": "App", "ceo": "CEO", "coo": "COO",
    "cto": "CTO", "db": "DB", "erp": "ERP", "hr": "HR", "it": "IT",
    "kpi": "KPI", "mis": "MIS", "mom": "MoM", "phd": "PhD", "roi": "ROI",
    "sop": "SOP", "vs": "VS",
}


# ============================================================================
# Font Utilities
# ============================================================================

def detect_font_path(font_name: str = "Arial") -> Optional[str]:
    """Auto-detect font path based on operating system."""
    system = platform.system()
    
    if system == "Windows":
        candidates = [
            Path(f"C:\\Windows\\Fonts\\{font_name.lower()}.ttf"),
            Path(f"C:\\Windows\\Fonts\\{font_name}.ttf"),
            Path(f"C:\\Windows\\Fonts\\{font_name}.ttc"),
        ]
    elif system == "Darwin":  # macOS
        candidates = [
            Path(f"/Library/Fonts/{font_name}.ttf"),
            Path(f"/Library/Fonts/{font_name}.ttc"),
            Path(f"/System/Library/Fonts/{font_name}.ttf"),
            Path(f"/System/Library/Fonts/Supplemental/{font_name}.ttf"),
        ]
    else:  # Linux
        candidates = [
            Path(f"/usr/share/fonts/truetype/msttcorefonts/{font_name}.ttf"),
            Path(f"/usr/share/fonts/{font_name}.ttf"),
            Path(f"/usr/share/fonts/truetype/{font_name}.ttf"),
        ]
    
    for path in candidates:
        if path.exists():
            return str(path)
    
    return None


def measure_text_width(text: str, size_pt: int, font_path: Optional[str] = None) -> float:
    """Measure the width of text in points using PIL ImageFont."""
    try:
        if font_path:
            font = ImageFont.truetype(font_path, int(size_pt))
        else:
            font = ImageFont.load_default()
        
        bbox = font.getbbox(text)
        return bbox[2] - bbox[0]
    except Exception:
        return len(text) * size_pt * 0.55


# ============================================================================
# Text Utilities
# ============================================================================

def is_sentence(text: str) -> bool:
    """Check if text appears to be a complete sentence."""
    text = text.strip()
    return text.endswith((".", "。", "!", "？", "?", ";", "；"))


def apply_title_case(text: str, small_words: Set[str], preserved_tokens: Dict[str, str]) -> str:
    """Apply English title case rules to text."""
    matches = list(ALPHA_RE.finditer(text))
    if not matches:
        return text
    
    first_start = matches[0].start()
    last_end = matches[-1].end()
    
    def capitalize_token(token: str) -> str:
        if token.isupper():
            return token
        lower = token.lower()
        if lower in preserved_tokens:
            return preserved_tokens[lower]
        return token.capitalize()
    
    def repl(m):
        token = m.group()
        lower = token.lower()
        is_first = (m.start() == first_start)
        is_last = (m.end() == last_end)
        
        if lower in ("s", "t") and m.start() > 0 and text[m.start() - 1] == "'":
            return lower
        
        if is_first or is_last or lower not in small_words:
            return capitalize_token(token)
        
        return lower
    
    return ALPHA_RE.sub(repl, text)


def capitalize_first(text: str) -> str:
    """Capitalize the first letter of text."""
    if not text:
        return text
    return text[0].upper() + text[1:]


def merge_paragraph_runs(paragraph) -> Tuple:
    """Merge all runs in a paragraph into the first run."""
    full_text = ""
    first_run = None
    base_size = None
    
    for run in paragraph.runs:
        full_text += run.text
        if first_run is None:
            first_run = run
        if base_size is None and run.font.size is not None:
            base_size = run.font.size.pt
    
    for run in paragraph.runs:
        if run is not first_run:
            run.text = ""
    
    return first_run, full_text, base_size


def wrap_lines(text: str, size_pt: int, max_width_pt: float, font_path: Optional[str] = None) -> list:
    """Wrap text into lines based on maximum width."""
    lines = []
    for paragraph in text.split("\n"):
        words = re.findall(r"\S+", paragraph)
        if not words:
            lines.append("")
            continue
        
        current_line = words[0]
        current_width = measure_text_width(current_line, size_pt, font_path)
        
        for word in words[1:]:
            word_width = measure_text_width(word, size_pt, font_path)
            space_width = measure_text_width(" ", size_pt, font_path)
            
            if current_width + space_width + word_width <= max_width_pt:
                current_line += " " + word
                current_width += space_width + word_width
            else:
                lines.append(current_line)
                current_line = word
                current_width = word_width
        
        lines.append(current_line)
    
    return lines


# ============================================================================
# Shape Utilities
# ============================================================================

def get_group_scale(shape) -> Tuple[float, float]:
    """Calculate scale factors for a group shape."""
    el = shape._element
    grpSpPr = el.find(qn("p:grpSpPr"))
    if grpSpPr is None:
        return (1.0, 1.0)
    
    xfrm = grpSpPr.find(qn("a:xfrm"))
    if xfrm is None:
        return (1.0, 1.0)
    
    ext = xfrm.find(qn("a:ext"))
    chExt = xfrm.find(qn("a:chExt"))
    
    if ext is None or chExt is None:
        return (1.0, 1.0)
    
    try:
        ext_cx = int(ext.get("cx", "0"))
        ext_cy = int(ext.get("cy", "0"))
        ch_cx = int(chExt.get("cx", "0"))
        ch_cy = int(chExt.get("cy", "0"))
    except (ValueError, TypeError):
        return (1.0, 1.0)
    
    if ch_cx <= 0 or ch_cy <= 0:
        return (1.0, 1.0)
    
    return (ext_cx / ch_cx, ext_cy / ch_cy)


def should_use_short_translation(shape_width_emu: int, shape_height_emu: int) -> bool:
    """Determine if a text box should use short/compact translations."""
    if shape_width_emu is None or shape_height_emu is None:
        return False
    
    width_in = shape_width_emu / 914400
    height_in = shape_height_emu / 914400
    
    return (width_in < 5.0) or (width_in < 6.0 and height_in < 1.0)


def fit_font_size(text: str, shape_width_emu: int, shape_height_emu: int,
                  base_size_pt: int, word_wrap: bool = True, margin_emu: int = 60000,
                  line_height_ratio: float = 1.35, min_size: int = 5,
                  font_path: Optional[str] = None) -> int:
    """Calculate the optimal font size to fit text within a shape."""
    if not text:
        return base_size_pt
    
    max_width_pt = (shape_width_emu - margin_emu) / 12700
    shape_height_pt = (shape_height_emu - margin_emu) / 12700
    
    if max_width_pt <= 0 or shape_height_pt <= 0:
        return base_size_pt
    
    def fits(size: int) -> bool:
        words = text.split()
        for word in words:
            if measure_text_width(word, size, font_path) > max_width_pt:
                return False
        
        if word_wrap:
            lines = wrap_lines(text, size, max_width_pt, font_path)
            line_height = size * line_height_ratio
            total_height = len(lines) * line_height
            return total_height <= shape_height_pt
        else:
            lines = text.split("\n")
            longest = max(lines, key=lambda x: measure_text_width(x, size, font_path))
            return measure_text_width(longest, size, font_path) <= max_width_pt
    
    if fits(base_size_pt):
        return base_size_pt
    
    new_size = base_size_pt
    while new_size > min_size:
        new_size -= 1
        if fits(new_size):
            return new_size
    
    return min_size


# ============================================================================
# SmartArt Processing
# ============================================================================

def set_xml_run_font_family(rPr, family: str = "Arial"):
    """Set latin/eastAsian/complexScript typefaces on an XML a:rPr element."""
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", family)


def process_smartart(shape, translate_func, set_font_func) -> bool:
    """Translate text inside a SmartArt graphicFrame."""
    if shape._element.tag != qn("p:graphicFrame"):
        return False
    
    graphicFrame = shape._element
    graphic = graphicFrame.find(qn("a:graphic"))
    if graphic is None:
        return False
    
    graphicData = graphic.find(qn("a:graphicData"))
    if graphicData is None:
        return False
    
    DGM_RELIDS = "{http://schemas.openxmlformats.org/drawingml/2006/diagram}relIds"
    relIds = graphicData.find(DGM_RELIDS)
    if relIds is None:
        return False
    
    dm_rid = relIds.get(qn("r:dm"))
    if not dm_rid:
        return False
    
    slide_part = shape.part
    try:
        data_part = slide_part.related_part(dm_rid)
    except Exception:
        return False
    
    root = etree.fromstring(data_part._blob)
    modified = False
    
    for t_el in root.iter(qn("a:t")):
        original = t_el.text or ""
        stripped = original.strip()
        
        if not stripped:
            continue
        
        translated = translate_func(stripped)
        
        if translated != stripped:
            t_el.text = translated
            modified = True
            
            r_el = t_el.getparent()
            if r_el is not None and r_el.tag == qn("a:r"):
                rPr = r_el.find(qn("a:rPr"))
                if rPr is None:
                    rPr = etree.SubElement(r_el, qn("a:rPr"))
                set_font_func(rPr)
    
    if modified:
        data_part._blob = etree.tostring(root, encoding="UTF-8", standalone=True, xml_declaration=True)
    
    return modified


# ============================================================================
# Main Translator Class
# ============================================================================

class PPTTranslator:
    """Universal PowerPoint translator."""
    
    def __init__(self, direction: str = "zh2en", target_font: str = "Arial",
                 terminology_file: Optional[str] = None,
                 title_case_enabled: bool = True):
        self.direction = direction
        self.target_font = target_font
        self.title_case_enabled = title_case_enabled
        self.terminology: Dict[str, str] = {}
        self.font_path = detect_font_path(target_font)
        
        if terminology_file:
            self.load_terminology(terminology_file)
    
    def load_terminology(self, file_path: str):
        """Load terminology from Excel or JSON file."""
        path = Path(file_path)
        
        if path.suffix.lower() in (".xlsx", ".xls"):
            self._load_from_excel(file_path)
        elif path.suffix.lower() == ".json":
            self._load_from_json(file_path)
        else:
            raise ValueError(f"Unsupported terminology file format: {path.suffix}")
    
    def _load_from_excel(self, file_path: str):
        """Load terminology from Excel file."""
        wb = openpyxl.load_workbook(file_path, data_only=True)
        ws = wb.active
        
        for row in ws.iter_rows(values_only=True):
            if row and len(row) >= 2 and row[0] and row[1]:
                source, target = str(row[0]).strip(), str(row[1]).strip()
                if source and target:
                    self.terminology[source] = target
        
        wb.close()
    
    def _load_from_json(self, file_path: str):
        """Load terminology from JSON file."""
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        
        if isinstance(data, dict):
            self.terminology = data
        elif isinstance(data, list):
            for item in data:
                if "source" in item and "target" in item:
                    self.terminology[item["source"]] = item["target"]
    
    def translate_text(self, text: str, prefer_short: bool = False) -> str:
        """Translate text using terminology and rules."""
        original = text
        stripped = text.strip()
        
        if not stripped:
            return text
        
        # Check terminology
        if original in self.terminology:
            translated = self.terminology[original]
        elif stripped in self.terminology:
            translated = self.terminology[stripped]
        else:
            translated = stripped
        
        # Apply formatting rules based on direction
        if self.direction == "zh2en":
            if is_sentence(translated):
                translated = capitalize_first(translated)
            elif self.title_case_enabled:
                translated = apply_title_case(translated, DEFAULT_SMALL_WORDS, DEFAULT_PRESERVED_TOKENS)
        
        return translated
    
    def set_run_font_family(self, run, family: Optional[str] = None):
        """Set latin, eastAsian, and complexScript typefaces for a run."""
        if family is None:
            family = self.target_font
        
        run.font.name = family
        rPr = run._r.get_or_add_rPr()
        
        for tag in ("a:latin", "a:ea", "a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = etree.SubElement(rPr, qn(tag))
            el.set("typeface", family)
    
    def process_text_frame(self, tf, shape_width_emu: int, shape_height_emu: int):
        """Process a text frame, translating and formatting all paragraphs."""
        full_text = "".join(para.text for para in tf.paragraphs).strip()
        prefer_short = should_use_short_translation(shape_width_emu, shape_height_emu)
        
        if full_text and full_text in self.terminology:
            translated = self.translate_text(full_text, prefer_short)
            
            first_para = tf.paragraphs[0]
            first_run = first_para.runs[0] if first_para.runs else first_para.add_run()
            
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ""
            
            first_run.text = translated
            self.set_run_font_family(first_run)
            
            base_size = first_run.font.size.pt if first_run.font.size else None
            if base_size:
                new_size = fit_font_size(translated, shape_width_emu, shape_height_emu,
                                        base_size, tf.word_wrap, font_path=self.font_path)
                first_run.font.size = Pt(new_size)
            
            return
        
        for para in tf.paragraphs:
            self._process_paragraph(para, shape_width_emu, shape_height_emu)
    
    def _process_paragraph(self, para, shape_width_emu: int, shape_height_emu: int):
        """Process a single paragraph."""
        run, original_text, base_size = merge_paragraph_runs(para)
        
        if run is None or not original_text.strip():
            return
        
        prefer_short = should_use_short_translation(shape_width_emu, shape_height_emu)
        translated = self.translate_text(original_text, prefer_short)
        run.text = translated
        self.set_run_font_family(run)
        
        if base_size:
            new_size = fit_font_size(translated, shape_width_emu, shape_height_emu,
                                    base_size, para.word_wrap, font_path=self.font_path)
            run.font.size = Pt(new_size)
    
    def process_shape(self, shape, scale_x: float = 1.0, scale_y: float = 1.0):
        """Process a single shape."""
        if hasattr(shape, "shapes"):
            gscale_x, gscale_y = get_group_scale(shape)
            new_scale_x = scale_x * gscale_x
            new_scale_y = scale_y * gscale_y
            
            for child in shape.shapes:
                self.process_shape(child, new_scale_x, new_scale_y)
            return
        
        if shape._element.tag == qn("p:graphicFrame"):
            process_smartart(
                shape,
                lambda text: self.translate_text(text),
                lambda rPr: set_xml_run_font_family(rPr, self.target_font),
            )
            return
        
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    self.process_text_frame(cell.text_frame,
                                          shape.width * scale_x,
                                          shape.height * scale_y)
            return
        
        if shape.has_text_frame:
            self.process_text_frame(shape.text_frame,
                                  shape.width * scale_x,
                                  shape.height * scale_y)
    
    def translate_ppt(self, input_path: str, output_path: str):
        """Translate an entire PowerPoint file."""
        prs = Presentation(input_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                self.process_shape(shape)
        
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        prs.save(output_path)
        print(f"Saved translated PPT to: {output_path}")


# ============================================================================
# CLI
# ============================================================================

def main():
    """Main CLI entry point."""
    parser = argparse.ArgumentParser(
        prog="ppt-translate",
        description="Translate PowerPoint presentations between Chinese and English",
    )
    
    parser.add_argument("input", type=str, help="Input PowerPoint file (.pptx)")
    parser.add_argument("-o", "--output", type=str, help="Output PowerPoint file")
    parser.add_argument("-d", "--direction", type=str, choices=["zh2en", "en2zh"],
                       default="zh2en", help="Translation direction (default: zh2en)")
    parser.add_argument("-t", "--terms", type=str, help="Terminology file (.xlsx or .json)")
    parser.add_argument("--font", type=str, default="Arial", help="Target font family")
    parser.add_argument("--no-title-case", action="store_true", help="Disable title case")
    
    args = parser.parse_args()
    
    try:
        # Validate input
        if not os.path.exists(args.input):
            print(f"Error: Input file not found: {args.input}", file=sys.stderr)
            return 1
        
        if not args.input.lower().endswith((".pptx", ".ppt")):
            print(f"Error: Input file is not a PowerPoint file", file=sys.stderr)
            return 1
        
        # Generate output path if not specified
        if not args.output:
            path = Path(args.input)
            suffix = "_EN" if args.direction == "zh2en" else "_ZH"
            args.output = str(path.parent / f"{path.stem}{suffix}{path.suffix}")
        
        # Create translator and run
        translator = PPTTranslator(
            direction=args.direction,
            target_font=args.font,
            terminology_file=args.terms,
            title_case_enabled=not args.no_title_case,
        )
        translator.translate_ppt(args.input, args.output)
        
        print(f"\nTranslation complete!")
        print(f"Direction: {args.direction}")
        print(f"Output: {args.output}")
        
        return 0
        
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
